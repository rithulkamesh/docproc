"""PowerPoint (.pptx) document loader."""

from pathlib import Path
from typing import Iterator

from docproc.doc.loaders.base import DocumentLoader, LoadedPage
from docproc.doc.regions import BoundingBox, Region, RegionType


class PPTXLoader(DocumentLoader):
    """Load PowerPoint files via python-pptx. Yields one LoadedPage per slide."""

    def load(self, path: Path) -> Iterator[LoadedPage]:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(path)
        for slide_num, slide in enumerate(prs.slides):
            regions = []
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    t = shape.text.strip()
                    regions.append(
                        Region(
                            region_type=RegionType.TEXT,
                            bbox=BoundingBox(x1=0, y1=0, x2=0, y2=0),
                            confidence=0.95,
                            content=t,
                            metadata={"page_num": slide_num, "slide": slide_num + 1},
                        )
                    )
                    texts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        row_text = " | ".join(cells)
                        if row_text:
                            regions.append(
                                Region(
                                    region_type=RegionType.TABLE,
                                    bbox=BoundingBox(x1=0, y1=0, x2=0, y2=0),
                                    confidence=0.95,
                                    content=row_text,
                                    metadata={"page_num": slide_num, "slide": slide_num + 1},
                                )
                            )
                            texts.append(row_text)
            images = []
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture
                    try:
                        img = shape.image
                        if img.blob:
                            images.append(img.blob)
                    except Exception:
                        pass
            yield LoadedPage(
                page_num=slide_num,
                text="\n\n".join(texts),
                regions=regions,
                raw_images=images,
            )

    def get_full_text(self, path: Path) -> str:
        from pptx import Presentation

        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
                if shape.has_table:
                    for row in shape.table.rows:
                        slide_texts.append(" | ".join(c.text.strip() for c in row.cells))
            if slide_texts:
                parts.append("\n".join(slide_texts))
        return "\n\n---\n\n".join(parts)
